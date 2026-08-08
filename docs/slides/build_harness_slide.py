#!/usr/bin/env python3
"""Build the "One goal, eight engines" slide.

Single source of truth: COPY (locked text) + build_scene() (geometry in CSS px,
1 inch = 100 px). Two renderers consume the same scene list:

  render_pptx()          -> harness-flow.pptx   (13.333 x 7.5 in, real shapes)
  render_html_preview()  -> preview.html        (1333 x 750 px, same coords)

Running `python3 build_harness_slide.py` produces both.

Layout: title + subtitle, then a slim header strip (spoken goal pill ->
Confirm gate -> elbow into engine 1), then a serpentine 3x3 grid:
row 1 engines 1-3 left-to-right, row 2 engines 4-6 right-to-left,
row 3 engine 7 -> Approve gate -> engine 8 (dashed border + "later" badge,
reached by a dashed connector: it wakes after the plan is made).
"""

import math
import os

BASE = os.path.dirname(os.path.abspath(__file__))

# ---------------------------------------------------------------- palette ----
BG          = '#ECEFF4'   # slide background
CARD        = '#FFFFFF'   # card fill
INK         = '#1A1E2C'   # primary ink
INK2        = '#5C6377'   # secondary ink
ACCENT      = '#3F6FE8'   # gates, connectors, numbers
ACCENT_TINT = '#E9EFFC'   # gate fill
GREEN       = '#12A150'   # "what it did" mark
GREEN_TINT  = '#E7F7EE'
HAIR        = '#E2E6F0'   # hairline / badge chip
BORDER      = '#868FA6'   # visible container border, used sparingly

# ------------------------------------------------------------- type scale ----
PT_TITLE = 26.0
PT_SUB   = 12.0
PT_NAME  = 12.5    # engine names (bold)
PT_PURP  = 10.5    # purpose line
PT_EX    = 9.5     # example lines (a step below the purpose)
PT_GNAME = 12.0    # gate names (bold)
PT_GLINE = 10.0    # gate one-liners
PT_NUM   = 11.0    # disc numbers
PT_BADGE = 9.5     # "later" badge
PT_QUOTE = 11.0    # spoken-goal pill

FONT_PPTX = 'Segoe UI'
FONT_CSS  = '"Segoe UI", "Liberation Sans", system-ui, -apple-system, sans-serif'

# ------------------------------------------------------------------- copy ----
TITLE    = 'ONE GOAL, EIGHT ENGINES'
SUBTITLE = 'What the Family Hub does with a sentence'
QUOTE    = '"Help the family eat healthier next week, and cut food waste."'

GATE_CONFIRM = ('Confirm', 'The family checks it heard right')
GATE_APPROVE = ('Approve', 'The plan goes to the board')

# (num, name, purpose, example 1, example 2)
ENGINES = [
    (1, 'Pre-check Engine', 'Make sure things work before it starts.',
        'The fridge is online and the account is signed in.',
        "If it weren't, it stops here — before writing a plan nobody could use."),
    (2, 'Capability Manager', 'Know what the fridge can do.',
        'Picks what this goal needs: food, recipes, activity, shopping.',
        "Leaves out the rest — it won't touch the door locks to plan dinner."),
    (3, 'Task Manager', 'Break the goal into jobs and follow them.',
        'Two jobs from one sentence: eat healthier, and waste less.',
        "Under them: find what's going off, choose meals, fix the shopping list."),
    (4, 'Grounding', 'Go and look at what is really there.',
        'Reads the fridge: spinach, yoghurt, bread and chicken go off in three days.',
        'Reads the family: a peanut allergy, no pork, low salt, and how much '
        'everyone has been moving.'),
    (5, 'Planner', 'Work out the week, day by day.',
        'Puts the spinach and chicken in the first dinners, so nothing goes in the bin.',
        'Turns down beef chilli — the family would rather have white meat — and says so.'),
    (6, 'Safety Policy Engine',
        'Decide what it may do alone, and what it must ask about.',
        'Pork never even reaches the list of options.',
        'Suggesting the shopping is fine; actually ordering it needs a yes.'),
    (7, 'Approval', 'Nothing real happens until someone says yes.',
        'Holds one shopping list: the few things missing for the week.',
        'Ticking off the spinach it used needs nobody — it just says it did.'),
    (8, 'Monitor & Adapt', 'Keep the plan right afterwards.',
        'Fish is delivered — it moves the fish forward before it spoils.',
        'Someone has a heavy training day — it raises the protein that evening, '
        'and says both reasons.'),
]
BADGE8 = 'Later, when a day passes'

# --------------------------------------------------------------- geometry ----
W, H   = 1333.0, 750.0
MARG   = 34.0
CARD_W = 400.0
GAP_X  = 32.0
XS     = [MARG, MARG + CARD_W + GAP_X, MARG + 2 * (CARD_W + GAP_X)]  # 34/466/898

TITLE_Y, SUB_Y = 12.0, 50.0
STRIP_Y, STRIP_H = 74.0, 42.0                 # goal pill + Confirm gate
ROW_Y  = [148.0, 351.0, 554.0]                # three card rows
CARD_H = 174.0

# card internals (offsets from card origin)
DISC_Y   = 12.0    # number disc / name band
PURP_Y   = 44.0    # purpose text (up to two lines)
CHIP1_Y  = 82.0    # first example chip
CHIP2_Y  = 127.0   # second example chip
CHIP_H   = 40.0

# ------------------------------------------------------- scene primitives ----

def run(t, size, color, bold=False, italic=False, spc=None):
    return dict(t=t, size=size, color=color, bold=bold, italic=italic, spc=spc)

def para(runs, align='left', leading=1.12):
    return dict(runs=runs, align=align, leading=leading)


def arrow_h(x_tail, x_tip, cy, shaft=4.0, head_l=11.0, head_w=14.0):
    d = 1.0 if x_tip > x_tail else -1.0
    xb = x_tip - d * head_l
    s, hw = shaft / 2.0, head_w / 2.0
    return [(x_tail, cy - s), (xb, cy - s), (xb, cy - hw), (x_tip, cy),
            (xb, cy + hw), (xb, cy + s), (x_tail, cy + s)]


def arrow_v(y_tail, y_tip, cx, shaft=4.0, head_l=11.0, head_w=14.0):
    d = 1.0 if y_tip > y_tail else -1.0
    yb = y_tip - d * head_l
    s, hw = shaft / 2.0, head_w / 2.0
    return [(cx - s, y_tail), (cx - s, yb), (cx - hw, yb), (cx, y_tip),
            (cx + hw, yb), (cx + s, yb), (cx + s, y_tail)]


def semicircle(cx, cy, r, steps=16):
    """Flat-bottomed semicircle (person shoulders)."""
    pts = []
    for i in range(steps + 1):
        a = math.pi * (1.0 - i / steps)
        pts.append((cx + r * math.cos(a), cy - r * math.sin(a)))
    return pts


def build_scene():
    S = []

    def rect(x, y, w, h, fill=None, line=None, lw=1.0, radius=0.0, dash=False):
        S.append(dict(kind='rect', x=x, y=y, w=w, h=h, fill=fill, line=line,
                      lw=lw, radius=radius, dash=dash))

    def ellipse(x, y, w, h, fill=None, line=None, lw=1.0):
        S.append(dict(kind='ellipse', x=x, y=y, w=w, h=h, fill=fill,
                      line=line, lw=lw))

    def poly(pts, fill):
        S.append(dict(kind='poly', pts=pts, fill=fill))

    def seg(x1, y1, x2, y2, color, lw=2.0, dash=False):
        S.append(dict(kind='line', x1=x1, y1=y1, x2=x2, y2=y2, color=color,
                      lw=lw, dash=dash))

    def text(x, y, w, h, paras, valign='top'):
        S.append(dict(kind='text', x=x, y=y, w=w, h=h, paras=paras,
                      valign=valign))

    def person(cx, cy):
        """Small human glyph: head + shoulders, accent."""
        ellipse(cx - 5.5, cy - 14.0, 11.0, 11.0, fill=ACCENT)
        poly(semicircle(cx, cy + 11.0, 10.0), ACCENT)

    def gate(x, y, w, h, name, line_txt):
        """A human decision: pill shape, accent tint — never a card."""
        rect(x, y, w, h, fill=ACCENT_TINT, line=ACCENT, lw=1.5, radius=h / 2.0)
        person(x + 32.0, y + h / 2.0)
        text(x + 54.0, y, w - 68.0, h,
             [para([run(name, PT_GNAME, INK, bold=True)]),
              para([run(line_txt, PT_GLINE, INK2)])], valign='middle')

    def engine_card(x, y, num, name, purpose, ex1, ex2, badge=None):
        dashed = badge is not None
        rect(x, y, CARD_W, CARD_H, fill=CARD,
             line=ACCENT if dashed else None, lw=1.4, radius=10.0, dash=dashed)
        ellipse(x + 16, y + DISC_Y, 26, 26, fill=ACCENT)
        text(x + 16, y + DISC_Y, 26, 26,
             [para([run(str(num), PT_NUM, '#FFFFFF', bold=True)],
                   align='center')], valign='middle')
        text(x + 52, y + DISC_Y, CARD_W - 68, 26,
             [para([run(name, PT_NAME, INK, bold=True)])], valign='middle')
        if badge:
            bw = 182.0
            rect(x + CARD_W - 14 - bw, y + 13, bw, 24, fill=HAIR, radius=12.0)
            text(x + CARD_W - 14 - bw, y + 13, bw, 24,
                 [para([run(badge, PT_BADGE, INK2)], align='center')],
                 valign='middle')
        text(x + 16, y + PURP_Y, CARD_W - 32, 34,
             [para([run(purpose, PT_PURP, INK2)], leading=1.15)])
        for cy0, ex in ((CHIP1_Y, ex1), (CHIP2_Y, ex2)):
            rect(x + 12, y + cy0, CARD_W - 24, CHIP_H, fill=GREEN_TINT,
                 radius=8.0)
            text(x + 22, y + cy0, 14, CHIP_H,
                 [para([run('✓', PT_EX, GREEN, bold=True)])], valign='middle')
            text(x + 40, y + cy0, CARD_W - 62, CHIP_H,
                 [para([run(ex, PT_EX, INK)], leading=1.15)], valign='middle')

    # background (drawn a hair wide: the true 16:9 slide is 13.333 in, i.e.
    # 1333.33 px; the HTML preview clips the overhang at exactly 1333 px)
    rect(0, 0, W + 1.0 / 3.0, H, fill=BG)

    # title block
    text(MARG, TITLE_Y, 900, 36,
         [para([run(TITLE, PT_TITLE, INK, bold=True, spc=20)])])
    text(MARG, SUB_Y, 1100, 20, [para([run(SUBTITLE, PT_SUB, INK2)])])

    # --- header strip: spoken goal -> Confirm gate -> elbow to engine 1 ----
    pill_w = 540.0
    rect(MARG, STRIP_Y, pill_w, STRIP_H, fill=CARD, line=BORDER, lw=1.2,
         radius=STRIP_H / 2.0)
    text(MARG + 18, STRIP_Y, pill_w - 36, STRIP_H,
         [para([run(QUOTE, PT_QUOTE, INK, italic=True)], align='center')],
         valign='middle')

    strip_cy = STRIP_Y + STRIP_H / 2.0
    gate_x, gate_w = MARG + pill_w + 44.0, 300.0          # 618 .. 918
    poly(arrow_h(MARG + pill_w + 8, gate_x - 8, strip_cy), ACCENT)
    gate(gate_x, STRIP_Y, gate_w, STRIP_H, *GATE_CONFIRM)

    # elbow: Confirm gate -> engine 1
    gcx = gate_x + gate_w / 2.0
    e1cx = XS[0] + CARD_W / 2.0
    ey = STRIP_Y + STRIP_H + 15.0                          # 131
    seg(gcx, STRIP_Y + STRIP_H + 1, gcx, ey, ACCENT, lw=2.2)
    seg(gcx + 1.1, ey, e1cx, ey, ACCENT, lw=2.2)
    poly(arrow_v(ey - 1, ROW_Y[0] - 3, e1cx), ACCENT)

    # --- row 1: engines 1-3, left to right --------------------------------
    cy1 = ROW_Y[0] + CARD_H / 2.0
    for i in range(3):
        engine_card(XS[i], ROW_Y[0], *ENGINES[i])
        if i < 2:
            gx = XS[i] + CARD_W
            poly(arrow_h(gx + 4, gx + GAP_X - 4, cy1), ACCENT)

    # down: engine 3 -> engine 4 (right column)
    poly(arrow_v(ROW_Y[0] + CARD_H + 3, ROW_Y[1] - 3, XS[2] + CARD_W / 2.0),
         ACCENT)

    # --- row 2: engines 4-6, right to left --------------------------------
    cy2 = ROW_Y[1] + CARD_H / 2.0
    for j, col in enumerate([2, 1, 0]):                    # engines 4, 5, 6
        engine_card(XS[col], ROW_Y[1], *ENGINES[3 + j])
        if col > 0:
            poly(arrow_h(XS[col] - 4, XS[col] - GAP_X + 4, cy2), ACCENT)

    # down: engine 6 -> engine 7 (left column)
    poly(arrow_v(ROW_Y[1] + CARD_H + 3, ROW_Y[2] - 3, XS[0] + CARD_W / 2.0),
         ACCENT)

    # --- row 3: engine 7 -> Approve gate -> (later) engine 8 --------------
    cy3 = ROW_Y[2] + CARD_H / 2.0
    engine_card(XS[0], ROW_Y[2], *ENGINES[6])

    ag_w, ag_h = 340.0, 90.0
    ag_x = XS[1] + (CARD_W - ag_w) / 2.0                   # 496
    ag_y = ROW_Y[2] + (CARD_H - ag_h) / 2.0                # 596
    poly(arrow_h(XS[0] + CARD_W + 4, ag_x - 4, cy3), ACCENT)
    gate(ag_x, ag_y, ag_w, ag_h, *GATE_APPROVE)

    # dashed connector: the board -> Monitor & Adapt, which wakes later
    seg(ag_x + ag_w + 4, cy3, XS[2] - 13, cy3, ACCENT, lw=2.2, dash=True)
    poly([(XS[2] - 13, cy3 - 7), (XS[2] - 2, cy3), (XS[2] - 13, cy3 + 7)],
         ACCENT)

    # engine 8: dashed border + badge = set apart from the live run
    engine_card(XS[2], ROW_Y[2], *ENGINES[7], badge=BADGE8)

    return S


# ---------------------------------------------------------- pptx renderer ----

def render_pptx(scene, path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE

    def E(px):
        return Emu(int(round(px * 9144)))          # 1 px = 1/100 in = 9144 EMU

    def C(hexstr):
        return RGBColor.from_string(hexstr.lstrip('#'))

    prs = Presentation()
    prs.slide_width = Emu(12192000)                # exactly 13.333 in (16:9)
    prs.slide_height = E(H)                        # exactly 7.5 in
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    def style_shape(shp, fill, line, lw, dash=False):
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = C(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = C(line)
            shp.line.width = Pt(lw)
            if dash:
                shp.line.dash_style = MSO_LINE.DASH
        else:
            shp.line.fill.background()
        try:
            shp.shadow.inherit = False
        except Exception:
            pass

    for el in scene:
        k = el['kind']
        if k == 'rect':
            if el['radius'] > 0:
                shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       E(el['x']), E(el['y']), E(el['w']), E(el['h']))
                shp.adjustments[0] = min(0.5, el['radius'] / min(el['w'], el['h']))
            else:
                shp = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            style_shape(shp, el['fill'], el['line'], el['lw'], el['dash'])
        elif k == 'ellipse':
            shp = shapes.add_shape(MSO_SHAPE.OVAL,
                                   E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            style_shape(shp, el['fill'], el['line'], el['lw'])
        elif k == 'poly':
            pts = el['pts']
            fb = shapes.build_freeform(E(pts[0][0]), E(pts[0][1]), scale=1.0)
            fb.add_line_segments([(E(x), E(y)) for x, y in pts[1:]], close=True)
            shp = fb.convert_to_shape()
            style_shape(shp, el['fill'], None, 0)
        elif k == 'line':
            conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        E(el['x1']), E(el['y1']),
                                        E(el['x2']), E(el['y2']))
            conn.line.color.rgb = C(el['color'])
            conn.line.width = Pt(el['lw'])
            if el['dash']:
                conn.line.dash_style = MSO_LINE.DASH
            try:
                conn.shadow.inherit = False
            except Exception:
                pass
        elif k == 'text':
            tb = shapes.add_textbox(E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = (MSO_ANCHOR.MIDDLE if el['valign'] == 'middle'
                                  else MSO_ANCHOR.TOP)
            for i, p in enumerate(el['paras']):
                pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pp.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                                'right': PP_ALIGN.RIGHT}[p['align']]
                pp.line_spacing = p['leading']
                for r in p['runs']:
                    rr = pp.add_run()
                    rr.text = r['t']
                    f = rr.font
                    f.name = FONT_PPTX
                    f.size = Pt(r['size'])
                    f.bold = r['bold']
                    f.italic = r['italic']
                    f.color.rgb = C(r['color'])
                    if r['spc'] is not None:
                        rr._r.get_or_add_rPr().set('spc', str(int(r['spc'])))

    prs.save(path)


# ---------------------------------------------------------- html renderer ----

def render_html_preview(scene, path):
    PX = 100.0 / 72.0                              # 1 pt = 1.389 px

    def css_border(el):
        if not el['line']:
            return ''
        style = 'dashed' if el['dash'] else 'solid'
        return f"border:{el['lw'] * PX:.2f}px {style} {el['line']};"

    out = []
    for el in scene:
        k = el['kind']
        if k == 'rect':
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;box-sizing:border-box;"
                f"background:{el['fill'] or 'transparent'};"
                f"border-radius:{el['radius']:.1f}px;{css_border(el)}\"></div>")
        elif k == 'ellipse':
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;box-sizing:border-box;"
                f"background:{el['fill'] or 'transparent'};border-radius:50%;"
                f"{css_border(dict(el, dash=False))}\"></div>")
        elif k == 'poly':
            pts = ' '.join(f"{x:.1f},{y:.1f}" for x, y in el['pts'])
            out.append(
                f"<svg style=\"position:absolute;left:0;top:0;pointer-events:none\" "
                f"width=\"{W:.0f}\" height=\"{H:.0f}\">"
                f"<polygon points=\"{pts}\" fill=\"{el['fill']}\"/></svg>")
        elif k == 'line':
            dash = ' stroke-dasharray="6 5"' if el['dash'] else ''
            out.append(
                f"<svg style=\"position:absolute;left:0;top:0;pointer-events:none\" "
                f"width=\"{W:.0f}\" height=\"{H:.0f}\">"
                f"<line x1=\"{el['x1']:.1f}\" y1=\"{el['y1']:.1f}\" "
                f"x2=\"{el['x2']:.1f}\" y2=\"{el['y2']:.1f}\" "
                f"stroke=\"{el['color']}\" stroke-width=\"{el['lw'] * PX:.2f}\"{dash}/>"
                f"</svg>")
        elif k == 'text':
            just = 'center' if el['valign'] == 'middle' else 'flex-start'
            paras_html = []
            for p in el['paras']:
                spans = []
                for r in p['runs']:
                    st = (f"font-size:{r['size'] * PX:.2f}px;color:{r['color']};"
                          f"font-weight:{700 if r['bold'] else 400};"
                          f"font-style:{'italic' if r['italic'] else 'normal'};")
                    if r['spc'] is not None:
                        st += f"letter-spacing:{r['spc'] / 100.0 * PX:.2f}px;"
                    spans.append(f"<span style=\"{st}\">{r['t']}</span>")
                paras_html.append(
                    f"<div style=\"text-align:{p['align']};"
                    f"line-height:{p['leading']:.2f}\">{''.join(spans)}</div>")
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;display:flex;"
                f"flex-direction:column;justify-content:{just};\">"
                f"{''.join(paras_html)}</div>")

    html = (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
        "<title>harness slide preview</title><style>"
        f"html,body{{margin:0;padding:0}}body{{font-family:{FONT_CSS};}}"
        "</style></head><body>"
        f"<div style=\"position:relative;width:{W:.0f}px;height:{H:.0f}px;"
        f"overflow:hidden\">{''.join(out)}</div>"
        "</body></html>")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html)


def check_bounds(scene):
    """No shape may leave the 1333.33 x 750 canvas."""
    bad = []
    for el in scene:
        k = el['kind']
        if k in ('rect', 'ellipse', 'text'):
            x0, y0, x1, y1 = el['x'], el['y'], el['x'] + el['w'], el['y'] + el['h']
        elif k == 'poly':
            xs = [p[0] for p in el['pts']]
            ys = [p[1] for p in el['pts']]
            x0, y0, x1, y1 = min(xs), min(ys), max(xs), max(ys)
        else:
            x0, x1 = sorted((el['x1'], el['x2']))
            y0, y1 = sorted((el['y1'], el['y2']))
        if x0 < 0 or y0 < 0 or x1 > W + 1.0 / 3.0 or y1 > H:
            bad.append((k, x0, y0, x1, y1))
    return bad


def main():
    scene = build_scene()
    bad = check_bounds(scene)
    if bad:
        for b in bad:
            print('OUT OF BOUNDS:', b)
        raise SystemExit(1)
    render_pptx(scene, os.path.join(BASE, 'harness-flow.pptx'))
    render_html_preview(scene, os.path.join(BASE, 'preview.html'))
    print('wrote harness-flow.pptx and preview.html')


if __name__ == '__main__':
    main()
