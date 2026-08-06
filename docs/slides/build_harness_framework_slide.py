#!/usr/bin/env python3
"""Build the "Nine harnesses" reference-table slide.

Single source of truth: COPY (locked text) + build_scene() (geometry in CSS px,
1 inch = 100 px). Two renderers consume the same scene list:

  render_pptx()          -> harness-framework-slide.pptx  (13.333 x 7.5 in)
  render_html_preview()  -> harness-framework-preview.html (1333 x 750 px)

Running `python3 build_harness_framework_slide.py` produces both.
"""

import os

BASE = os.path.dirname(os.path.abspath(__file__))

# ---------------------------------------------------------------- palette ----
BG        = '#ECEFF4'   # slide background
CARD      = '#FFFFFF'   # table card / odd rows
TINT      = '#F7F9FC'   # even-row tint, barely there
NAVY      = '#1A1E2C'   # header band + primary ink
INK       = '#1A1E2C'
INK2      = '#5C6377'   # secondary ink (For / Does columns)
ACCENT    = '#3F6FE8'   # number discs
DOT       = '#95B0F2'   # hub-item dot (accent at ~55% on white)
DASH      = '#AEB8CC'   # does-item dash
HAIR      = '#E2E6F0'   # hairlines between rows

# ------------------------------------------------------------- type scale ----
PT_TITLE = 22.0
PT_SUB   = 12.0
PT_HDR   = 9.0     # column headers
PT_NAME  = 10.5    # harness names (bold)
PT_BODY  = 10.0    # everything in the three text columns
PT_NUM   = 10.5    # disc numbers

FONT_PPTX = 'Segoe UI'
FONT_CSS  = '"Segoe UI", "Liberation Sans", system-ui, -apple-system, sans-serif'

# ------------------------------------------------------------------- copy ----
TITLE    = 'THE NINE HARNESSES, AND WHAT EACH ONE DOES'
SUBTITLE = ('Every part of the agent that runs on the fridge — '
            'all nine are built and working')

HDR_NUM, HDR_NAME, HDR_FOR, HDR_DOES, HDR_HUB = (
    '#', 'HARNESS', 'WHAT IT IS FOR', 'WHAT IT DOES', 'ON A FAMILY HUB')

# (num, name lines, what it is for, what it does, on a family hub)
ROWS = [
    (1, ['Pre-check Engine'],
     'Make sure things work before it starts.',
     ['Checks what it needs is switched on',
      'Checks again just before it acts',
      'Stops with a fix you can do'],
     ['Signed in to the account',
      'Internet and smart home working',
      'Oven or fridge switched on']),
    (2, ['Capability Manager'],
     'Know what the fridge can do.',
     ['Finds everything the fridge can do',
      'Keeps looking apart from doing',
      'Hides what it is not allowed to do'],
     ['Food in the fridge, recipes',
      'Calendar, reminders, who lives here',
      'Shopping, deliveries, appliances']),
    (3, ['Task Manager'],
     'Break a goal into jobs and follow them.',
     ['Splits one goal into smaller jobs',
      'Follows each job to the end',
      'Pauses, retries or drops a job'],
     ["Plan the week's meals",
      'Get the home ready before a trip',
      'Plan a birthday or a dinner']),
    (4, ['Context & Grounding', 'Engine'],
     'Look at what is really going on at home.',
     ['Reads what is true right now',
      'Only looks, never changes anything',
      'Knows what day it is'],
     ['What is in the fridge, what is going off',
      'What the family has on this week',
      "What the household will and won't eat"]),
    (5, ['Goal Planner'],
     'Work out the plan, day by day.',
     ['Weighs the choices against house rules',
      'Puts a day against every step',
      'Says what it turned down, and why'],
     ['Seven dinners, each with a reason',
      "Jobs before you go, jobs for when you're back",
      "A dish dropped because you'd rather not"]),
    (6, ['Safety Policy Engine'],
     'Decide what it may do alone, and what it must ask about.',
     ['Sorts every action into four levels',
      'Holds the house rules on every action',
      'Drops unsafe choices before they are offered'],
     ['Just do it — tick off food that was used',
      'Tell you — add to the shopping list',
      'Ask first — buy something, stop a delivery',
      'Never — there is no way to allow it']),
    (7, ['Approval Manager'],
     'Nothing real happens until someone says yes.',
     ['Holds back anything that changes things',
      'Puts it all in one list to agree to',
      'Only does what was agreed'],
     ["Say yes to the week's shopping",
      'Say yes to pausing deliveries',
      'Say no to one thing, keep the rest']),
    (8, ['Product API Adapter'],
     'The one place the agent talks to the fridge.',
     ['Everything it reads or changes goes through here',
      'Hides the differences between products',
      'Change the product, keep the agent'],
     ['Food and recipes',
      'Calendar and family',
      'Appliances, security, deliveries']),
    (9, ['Monitor & Adapt'],
     'Keep the plan right after it is agreed.',
     ['Keeps watching the home',
      'Works out whether a change matters',
      'Only redoes the part that changed'],
     ['Food arrives, or is about to go off',
      'More guests than planned',
      'Another plan takes some days away']),
]

# --------------------------------------------------------------- geometry ----
W, H  = 1333.0, 750.0
MARG  = 24.0                      # side margins; card spans MARG .. W - MARG
CARD_X, CARD_W = MARG, W - 2 * MARG          # 24 .. 1309, width 1285

TITLE_Y, SUB_Y = 16.0, 52.0

HDR_Y, HDR_H = 78.0, 30.0                    # navy band 78 .. 108
ROWS_Y = HDR_Y + HDR_H                       # rows start at 108

ROW_H      = 66.0                            # normal row (3 list items)
ROW_H_TALL = 90.0                            # row 6 (4 hub items)
LINE_H     = 18.0                            # one list line
PAD        = 6.0                             # top/bottom padding, normal row
PAD_TALL   = 9.0                             # top/bottom padding, tall row

# columns: [x, width]; text gets a left pad inside each
COL_NUM  = (CARD_X,  40.0)                   #   24 ..   64
COL_NAME = (64.0,   172.0)                   #   64 ..  236
COL_FOR  = (236.0,  384.0)                   #  236 ..  620
COL_DOES = (620.0,  354.0)                   #  620 ..  974
COL_HUB  = (974.0,  335.0)                   #  974 .. 1309

NAME_TX   = COL_NAME[0] + 10.0               # 74
FOR_TX    = COL_FOR[0] + 12.0                # 248
DOES_MX   = COL_DOES[0] + 12.0               # 632  dash marker
DOES_TX   = DOES_MX + 15.0                   # 647
HUB_MX    = COL_HUB[0] + 12.0                # 986  dot marker
HUB_TX    = HUB_MX + 14.0                    # 1000
NAME_TW   = COL_NAME[0] + COL_NAME[1] - 8.0 - NAME_TX     # 154
FOR_TW    = COL_FOR[0] + COL_FOR[1] - 10.0 - FOR_TX       # 362
DOES_TW   = COL_DOES[0] + COL_DOES[1] - 10.0 - DOES_TX    # 317
HUB_TW    = COL_HUB[0] + COL_HUB[1] - 14.0 - HUB_TX       # 295

DISC = 26.0                                  # number disc diameter

# ------------------------------------------------------- scene primitives ----

def run(t, size, color, bold=False, italic=False, spc=None):
    return dict(t=t, size=size, color=color, bold=bold, italic=italic, spc=spc)

def para(runs, align='left', leading=1.0):
    return dict(runs=runs, align=align, leading=leading)


def build_scene():
    S = []

    def rect(x, y, w, h, fill=None, line=None, lw=1.0, radius=0.0, dash=False):
        S.append(dict(kind='rect', x=x, y=y, w=w, h=h, fill=fill, line=line,
                      lw=lw, radius=radius, dash=dash))

    def ellipse(x, y, w, h, fill=None, line=None, lw=1.0):
        S.append(dict(kind='ellipse', x=x, y=y, w=w, h=h, fill=fill,
                      line=line, lw=lw))

    def seg(x1, y1, x2, y2, color, lw=1.0, dash=False):
        S.append(dict(kind='line', x1=x1, y1=y1, x2=x2, y2=y2, color=color,
                      lw=lw, dash=dash))

    def text(x, y, w, h, paras, valign='top'):
        S.append(dict(kind='text', x=x, y=y, w=w, h=h, paras=paras,
                      valign=valign))

    # background (drawn a hair wide: the true 16:9 slide is 13.333 in, i.e.
    # 1333.33 px; the HTML preview clips the overhang at exactly 1333 px)
    rect(0, 0, W + 1.0 / 3.0, H, fill=BG)

    # title block
    text(CARD_X + 2, TITLE_Y, 1100, 34,
         [para([run(TITLE, PT_TITLE, INK, bold=True, spc=20)])])
    text(CARD_X + 2, SUB_Y, 1100, 20,
         [para([run(SUBTITLE, PT_SUB, INK2)])])

    # row grid: y positions and heights
    row_geo = []
    y = ROWS_Y
    for num, *_ in ROWS:
        h = ROW_H_TALL if num == 6 else ROW_H
        row_geo.append((y, h))
        y += h
    rows_bottom = y                                  # 726

    # card: navy header band (rounded top), white row area (rounded bottom)
    rect(CARD_X, HDR_Y, CARD_W, HDR_H + 10, fill=NAVY, radius=8.0)
    rect(CARD_X, ROWS_Y, CARD_W, rows_bottom - ROWS_Y, fill=CARD, radius=8.0)
    rect(CARD_X, ROWS_Y, 16, 16, fill=CARD)                       # square off
    rect(CARD_X + CARD_W - 16, ROWS_Y, 16, 16, fill=CARD)         # top corners

    # even-row tint (rows 2, 4, 6, 8)
    for i, (ry, rh) in enumerate(row_geo):
        if (i + 1) % 2 == 0:
            rect(CARD_X, ry, CARD_W, rh, fill=TINT)

    # hairlines between rows
    for ry, _ in row_geo[1:]:
        seg(CARD_X, ry, CARD_X + CARD_W, ry, HAIR, lw=1.0)

    # column headers (white small caps on navy)
    def hdr(t, x, w, align='left'):
        text(x, HDR_Y, w, HDR_H,
             [para([run(t, PT_HDR, '#FFFFFF', bold=True, spc=110)],
                   align=align)], valign='middle')

    hdr(HDR_NUM,  COL_NUM[0], COL_NUM[1], align='center')
    hdr(HDR_NAME, NAME_TX, COL_NAME[1] - 10)
    hdr(HDR_FOR,  FOR_TX, COL_FOR[1] - 12)
    hdr(HDR_DOES, DOES_MX, COL_DOES[1] - 12)
    hdr(HDR_HUB,  HUB_MX, COL_HUB[1] - 12)

    # rows
    for (num, name_lines, for_txt, does, hub), (ry, rh) in zip(ROWS, row_geo):
        pad = PAD_TALL if num == 6 else PAD
        top = ry + pad

        # number disc, centred on the first text line
        dy = top + LINE_H / 2.0 - DISC / 2.0
        dx = COL_NUM[0] + COL_NUM[1] / 2.0 - DISC / 2.0
        ellipse(dx, dy, DISC, DISC, fill=ACCENT)
        text(dx, dy, DISC, DISC,
             [para([run(str(num), PT_NUM, '#FFFFFF', bold=True)],
                   align='center')], valign='middle')

        # harness name (one or two lines, aligned with the first list line)
        text(NAME_TX, top, NAME_TW, LINE_H * len(name_lines),
             [para([run(ln, PT_NAME, INK, bold=True)], leading=1.25)
              for ln in name_lines], valign='middle')

        # what it is for — a single line
        text(FOR_TX, top, FOR_TW, LINE_H,
             [para([run(for_txt, PT_BODY, INK2)])], valign='middle')

        # what it does — short quiet dashes
        for i, item in enumerate(does):
            ly = top + i * LINE_H
            rect(DOES_MX, ly + LINE_H / 2.0 - 0.75, 8.0, 1.5, fill=DASH)
            text(DOES_TX, ly, DOES_TW, LINE_H,
                 [para([run(item, PT_BODY, INK2)])], valign='middle')

        # on a family hub — small accent dots
        for i, item in enumerate(hub):
            ly = top + i * LINE_H
            ellipse(HUB_MX, ly + LINE_H / 2.0 - 2.5, 5.0, 5.0, fill=DOT)
            text(HUB_TX, ly, HUB_TW, LINE_H,
                 [para([run(item, PT_BODY, INK)])], valign='middle')

    return S


# ---------------------------------------------------------- pptx renderer ----

def render_pptx(scene, path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE

    def E(px):
        return Emu(int(round(px * 9144)))          # 1 px = 1/100 in = 9144 EMU

    def C(hexstr):
        return RGBColor.from_string(hexstr.lstrip('#'))

    prs = Presentation()
    prs.slide_width = Emu(12192000)            # exactly 13.333 in (16:9)
    prs.slide_height = E(H)                    # exactly 7.5 in
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    def style_shape(shp, fill, line, lw, dash=False):
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = C(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = C(line)
            shp.line.width = Pt(lw)
            if dash:
                shp.line.dash_style = MSO_LINE.DASH
        else:
            shp.line.fill.background()
        try:
            shp.shadow.inherit = False
        except Exception:
            pass

    for el in scene:
        k = el['kind']
        if k == 'rect':
            if el['radius'] > 0:
                shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       E(el['x']), E(el['y']), E(el['w']), E(el['h']))
                shp.adjustments[0] = min(0.5, el['radius'] / min(el['w'], el['h']))
            else:
                shp = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            style_shape(shp, el['fill'], el['line'], el['lw'], el['dash'])
        elif k == 'ellipse':
            shp = shapes.add_shape(MSO_SHAPE.OVAL,
                                   E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            style_shape(shp, el['fill'], el['line'], el['lw'])
        elif k == 'line':
            conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        E(el['x1']), E(el['y1']),
                                        E(el['x2']), E(el['y2']))
            conn.line.color.rgb = C(el['color'])
            conn.line.width = Pt(el['lw'])
            if el['dash']:
                conn.line.dash_style = MSO_LINE.DASH
            try:
                conn.shadow.inherit = False
            except Exception:
                pass
        elif k == 'text':
            tb = shapes.add_textbox(E(el['x']), E(el['y']), E(el['w']), E(el['h']))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = (MSO_ANCHOR.MIDDLE if el['valign'] == 'middle'
                                  else MSO_ANCHOR.TOP)
            for i, p in enumerate(el['paras']):
                pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pp.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                                'right': PP_ALIGN.RIGHT}[p['align']]
                pp.line_spacing = p['leading']
                for r in p['runs']:
                    rr = pp.add_run()
                    rr.text = r['t']
                    f = rr.font
                    f.name = FONT_PPTX
                    f.size = Pt(r['size'])
                    f.bold = r['bold']
                    f.italic = r['italic']
                    f.color.rgb = C(r['color'])
                    if r['spc'] is not None:
                        rr._r.get_or_add_rPr().set('spc', str(int(r['spc'])))

    prs.save(path)


# ---------------------------------------------------------- html renderer ----

def render_html_preview(scene, path):
    PX = 100.0 / 72.0                              # 1 pt = 1.389 px

    def css_border(el):
        if not el['line']:
            return ''
        style = 'dashed' if el['dash'] else 'solid'
        return f"border:{el['lw'] * PX:.2f}px {style} {el['line']};"

    out = []
    for el in scene:
        k = el['kind']
        if k == 'rect':
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;box-sizing:border-box;"
                f"background:{el['fill'] or 'transparent'};"
                f"border-radius:{el['radius']:.1f}px;{css_border(el)}\"></div>")
        elif k == 'ellipse':
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;box-sizing:border-box;"
                f"background:{el['fill'] or 'transparent'};border-radius:50%;"
                f"{css_border(dict(el, dash=False))}\"></div>")
        elif k == 'line':
            dash = ' stroke-dasharray="6 5"' if el['dash'] else ''
            out.append(
                f"<svg style=\"position:absolute;left:0;top:0;pointer-events:none\" "
                f"width=\"{W:.0f}\" height=\"{H:.0f}\">"
                f"<line x1=\"{el['x1']:.1f}\" y1=\"{el['y1']:.1f}\" "
                f"x2=\"{el['x2']:.1f}\" y2=\"{el['y2']:.1f}\" "
                f"stroke=\"{el['color']}\" stroke-width=\"{el['lw'] * PX:.2f}\"{dash}/>"
                f"</svg>")
        elif k == 'text':
            just = 'center' if el['valign'] == 'middle' else 'flex-start'
            paras_html = []
            for p in el['paras']:
                spans = []
                for r in p['runs']:
                    st = (f"font-size:{r['size'] * PX:.2f}px;color:{r['color']};"
                          f"font-weight:{700 if r['bold'] else 400};"
                          f"font-style:{'italic' if r['italic'] else 'normal'};")
                    if r['spc'] is not None:
                        st += f"letter-spacing:{r['spc'] / 100.0 * PX:.2f}px;"
                    spans.append(f"<span style=\"{st}\">{r['t']}</span>")
                paras_html.append(
                    f"<div style=\"text-align:{p['align']};"
                    f"line-height:{p['leading']:.2f}\">{''.join(spans)}</div>")
            out.append(
                f"<div style=\"position:absolute;left:{el['x']:.1f}px;top:{el['y']:.1f}px;"
                f"width:{el['w']:.1f}px;height:{el['h']:.1f}px;display:flex;"
                f"flex-direction:column;justify-content:{just};\">"
                f"{''.join(paras_html)}</div>")

    html = (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\">"
        "<title>harness framework slide preview</title><style>"
        f"html,body{{margin:0;padding:0}}body{{font-family:{FONT_CSS};}}"
        "</style></head><body>"
        f"<div style=\"position:relative;width:{W:.0f}px;height:{H:.0f}px;"
        f"overflow:hidden\">{''.join(out)}</div>"
        "</body></html>")
    with open(path, 'w', encoding='utf-8') as f:
        f.write(html)


def main():
    scene = build_scene()
    render_pptx(scene, os.path.join(BASE, 'harness-framework-slide.pptx'))
    render_html_preview(scene, os.path.join(BASE, 'harness-framework-preview.html'))
    print('wrote harness-framework-slide.pptx and harness-framework-preview.html')


if __name__ == '__main__':
    main()
